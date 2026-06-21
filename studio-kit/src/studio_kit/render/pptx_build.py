"""把图片序列组装成横版 16:9 PPT（每图铺满一页，文案进演讲者备注）。

输入沿用 ArchVideoDoc（与视频共用同一份 script.json）：每段一页幻灯片，
seg.image 这张图按比例缩放居中铺满 16:9 页面，seg.narration 写入该页备注。
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from studio_kit.core.contracts import ArchVideoDoc
from studio_kit.core.logging import get_logger

logger = get_logger(__name__)

# 16:9 宽屏：13.333in × 7.5in（EMU：914400 / inch）
SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000


def build_pptx(
    doc: ArchVideoDoc,
    work_dir: Path,
    out_pptx: Path,
    *,
    notes: bool = True,
) -> Path:
    """生成横版 16:9 PPT。每段一页：图片按比例缩放居中，文案写入备注。

    seg.image 相对 work_dir 解析；缺图即 raise FileNotFoundError（不静默兜底）。
    """
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank_layout = prs.slide_layouts[6]  # 空白版式

    for seg in doc.segments:
        img = (work_dir / seg.image).resolve()
        if not img.exists():
            raise FileNotFoundError(f"第 {seg.index} 页引用的图片不存在：{img}")

        slide = prs.slides.add_slide(blank_layout)

        # 先以原始尺寸放入，再按比例缩放至「适配页面」并居中
        pic = slide.shapes.add_picture(str(img), 0, 0)
        scale = min(SLIDE_W_EMU / pic.width, SLIDE_H_EMU / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int((SLIDE_W_EMU - pic.width) / 2)
        pic.top = int((SLIDE_H_EMU - pic.height) / 2)

        if notes and seg.narration.strip():
            slide.notes_slide.notes_text_frame.text = seg.narration

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    logger.info("PPT 已生成（%d 页）：%s", len(doc.segments), out_pptx)
    return out_pptx
