from pathlib import Path

from PIL import Image
from pptx import Presentation

from studio_kit.core.contracts import ArchVideoDoc
from studio_kit.render.pptx_build import build_pptx, SLIDE_W_EMU, SLIDE_H_EMU


def _make_png(path: Path, w: int, h: int) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (w, h), (18, 20, 28)).save(path)


def _doc():
    return ArchVideoDoc.model_validate({
        "slug": "demo", "run_id": "r1", "title": "演示",
        "segments": [
            {"index": 0, "narration": "第一页讲解", "image": "images/a.png"},
            {"index": 1, "narration": "第二页讲解", "image": "images/b.png"},
        ],
    })


def test_build_pptx_slides_and_notes(tmp_path: Path):
    _make_png(tmp_path / "images" / "a.png", 1920, 1080)
    _make_png(tmp_path / "images" / "b.png", 1920, 1080)
    out = tmp_path / "demo.pptx"
    build_pptx(_doc(), tmp_path, out)
    assert out.exists()

    prs = Presentation(str(out))
    # 16:9 横版
    assert prs.slide_width == SLIDE_W_EMU
    assert prs.slide_height == SLIDE_H_EMU
    # 每段一页
    assert len(prs.slides) == 2
    # 文案进备注
    notes = [s.notes_slide.notes_text_frame.text for s in prs.slides]
    assert notes[0] == "第一页讲解"
    assert notes[1] == "第二页讲解"
    # 每页有一张图
    for s in prs.slides:
        assert any(sh.shape_type == 13 for sh in s.shapes)  # 13 = PICTURE


def test_build_pptx_fallback_when_locked(tmp_path: Path, monkeypatch):
    # 目标被占用（PermissionError）→ 自动另存 <名>-2.pptx，返回实际路径
    _make_png(tmp_path / "images" / "a.png", 1920, 1080)
    _make_png(tmp_path / "images" / "b.png", 1920, 1080)
    out = tmp_path / "demo.pptx"

    import pptx.presentation as _pp
    real_save = _pp.Presentation.save
    locked = {str(out)}

    def fake_save(self, path):
        if str(path) in locked:
            raise PermissionError(13, "in use")
        return real_save(self, path)

    monkeypatch.setattr(_pp.Presentation, "save", fake_save)
    saved = build_pptx(_doc(), tmp_path, out)
    assert saved == tmp_path / "demo-2.pptx"
    assert saved.exists()


def test_build_pptx_missing_image_raises(tmp_path: Path):
    import pytest
    # 只建第一张，第二张缺失 → raise
    _make_png(tmp_path / "images" / "a.png", 1920, 1080)
    with pytest.raises(FileNotFoundError):
        build_pptx(_doc(), tmp_path, tmp_path / "x.pptx")
